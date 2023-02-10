import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Pt
from pptx.util import Inches
from PIL import Image
import os

def create_pptx(titles, subtitles, image_paths, logo_path):
    # Create a new PowerPoint presentation
    prs = Presentation()
    # set width and height to 16 and 9 inches.
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)        

    # Loop over the titles, subtitles, and image paths
    for title, subtitle, image_path, logo_path in zip(titles, subtitles, image_paths, logo_path):
        # Add a slide to the presentation
        slide = prs.slides.add_slide(prs.slide_layouts[0])

        # Add the title to the slide
        title_text_frame = slide.shapes.add_textbox(Inches(2), Inches(0), Inches(10), Inches(1))
        title_text_frame.text = title
        title_text_frame.text_frame.paragraphs[0].font.size = Pt(44)


        # Add the subtitle to the slide
        subtitle_text_frame = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(10), Inches(1))
        subtitle_text_frame.text = subtitle
        subtitle_text_frame.text_frame.paragraphs[0].font.size = Pt(32)
        
        # Add the image to the slide
        left = Inches(2)
        top = Inches(2)
        height = Inches(5)
        pic = slide.shapes.add_picture(image_path, left, top, height=height)

        # Add the logo to the slide
        left = Inches(2)
        top = Inches(2)
        height = Inches(0.5)
        pic = slide.shapes.add_picture(logo_path, left, top, height=height)

    # Save the presentation
    prs.save("assignment_sample_Task.pptx")
    os.startfile("assignment_sample_Task.pptx")

# Slide Information
titles = ["Simple Title 1", "Simple Title 2", "Simple Title 3", "Simple Title 4", " Simple Title 5"]
subtitles = ["Simple Subtitle 1", "Simple Subtitle 2", "Simple Subtitle 3", "Simple Subtitle 4", "Simple Subtitle 5"]
image_paths = ["image1.jpg", "image2.jpg", "image3.jpg", "image4.jpg","image5.jpg" ]
logo_path=["nike_black.png", "nike_black.png", "nike_black.png", "nike_black.png", "nike_black.png"]

create_pptx(titles, subtitles, image_paths, logo_path)
